import argparse
import json
import os
import re
import tempfile
from collections import defaultdict
from typing import List, Dict, Any, Tuple, Optional, Set

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image
import pytesseract
import google.generativeai as genai


NUMBER_RE = re.compile(r'(?P<number>(?:\$|€|£)?\s?\d{1,3}(?:[,.\d{3}]+)?(?:\.\d+)?\s*(?:[kKmM])?(?:\s*(?:m|min|mins|minutes))?\s*%?)')
YEAR_RE = re.compile(r'\b(19|20)\d{2}\b')
PERCENT_RE = re.compile(r'(?P<pct>\d+(?:\.\d+)?)\s*%')

def extract_text_from_shape(shape) -> str:
    texts = []
    if not shape:
        return ""
    if hasattr(shape, "text") and shape.text:
        texts.append(shape.text)
    if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
        try:
            table = shape.table
            for r in range(len(table.rows)):
                row_texts = []
                for c in range(len(table.columns)):
                    cell = table.cell(r, c)
                    row_texts.append(cell.text)
                texts.append(" | ".join(row_texts))
        except Exception:
            pass
    return "\n".join([t for t in texts if t])

def extract_notes_text(slide) -> str:
    note_texts = []
    if slide.has_notes_slide:
        notes_slide = slide.notes_slide
        for shape in notes_slide.shapes:
            if hasattr(shape, "text") and shape.text:
                note_texts.append(shape.text)
    return "\n".join(note_texts)

def extract_images_from_pptx(prs: Presentation, out_dir: str) -> Dict[int, List[str]]:
    os.makedirs(out_dir, exist_ok=True)
    slide_images = defaultdict(list)
    for i, slide in enumerate(prs.slides, start=1):
        for shape_idx, shape in enumerate(slide.shapes):
            try:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    img = shape.image
                    ext = img.ext
                    fname = os.path.join(out_dir, f"slide_{i}_img_{shape_idx + 1}.{ext}")
                    with open(fname, "wb") as fh:
                        fh.write(img.blob)
                    slide_images[i].append(fname)
            except Exception:
                continue
    return slide_images

def ocr_image_file(image_path: str) -> str:
    try:
        img = Image.open(image_path)
        text = pytesseract.image_to_string(img)
        return text.strip()
    except Exception:
        return ""

def collect_slide_text(prs: Presentation, slide_idx: int) -> str:
    slide = prs.slides[slide_idx - 1]
    all_texts = []
    for shape in slide.shapes:
        try:
            t = extract_text_from_shape(shape)
            if t:
                all_texts.append(t)
        except Exception:
            continue
    notes = extract_notes_text(slide)
    if notes:
        all_texts.append("NOTES: " + notes)
    return "\n".join(all_texts)

def find_numbers(text: str) -> List[str]:
    return [m.group("number").strip() for m in NUMBER_RE.finditer(text)]

def find_years(text: str) -> List[int]:
    return [int(m.group(0)) for m in YEAR_RE.finditer(text)]

def analyze_pptx_in_memory(input_path: str, img_out_dir: str, ocr: bool = True) -> List[Dict[str,Any]]:
    prs = Presentation(input_path)
    slide_images = extract_images_from_pptx(prs, img_out_dir)
    slides_data = []
    for idx in range(1, len(prs.slides) + 1):
        raw = collect_slide_text(prs, idx)
        ocr_texts = []
        for img_path in slide_images.get(idx, []):
            if ocr:
                ocr_t = ocr_image_file(img_path)
                if ocr_t:
                    ocr_texts.append(ocr_t)
        combined_text = raw
        if ocr_texts:
            combined_text += "\n\nOCR_TEXT:\n" + "\n".join(ocr_texts)
        numbers = find_numbers(combined_text)
        years = find_years(combined_text)
        slides_data.append({
            "slide_num": idx,
            "text": combined_text,
            "raw_text": raw,
            "image_paths": slide_images.get(idx, []),
            "ocr_texts": ocr_texts,
            "numbers": numbers,
            "years": years
        })
    return slides_data


def normalize_number_token(tok: str) -> Tuple[Optional[float], str]:
    s = tok.strip().replace('\u200b','')
    s = re.sub(r'\s+', ' ', s)
    if '%' in s:
        s2 = s.replace('%','').replace(',','').strip()
        try:
            return float(s2), '%'
        except:
            return None, '%'
    if s and s[0] in ('$','€','£'):
        cur = s[0]
        s2 = s[1:].replace(',','').strip()
        mul = 1.0
        if s2.lower().endswith('k'):
            mul = 1e3; s2 = s2[:-1]
        elif s2.lower().endswith('m'):
            mul = 1e6; s2 = s2[:-1]
        try:
            return float(s2) * mul, f'currency:{cur}'
        except:
            return None, f'currency:{cur}'
    if re.search(r'\b(min|mins|minutes|m)\b', s, flags=re.IGNORECASE):
        s2 = re.sub(r'[^\d\.]', '', s)
        try:
            return float(s2), 'minutes'
        except:
            return None, 'minutes'
    s2 = s.replace(',','')
    mul = 1.0
    if s2.lower().endswith('k'):
        mul = 1e3; s2 = s2[:-1]
    elif s2.lower().endswith('m'):
        mul = 1e6; s2 = s2[:-1]
    try:
        return float(s2) * mul, ''
    except:
        return None, ''

def get_context(text: str, token: str, window_words: int = 6) -> str:
    words = re.findall(r'\w+|[%$€£]|[^\s\w]', text)
    token_l = token.strip()
    for i, w in enumerate(words):
        if token_l.lower() in w.lower() or w.lower() in token_l.lower():
            start = max(0, i - window_words)
            end = min(len(words), i + window_words + 1)
            return " ".join(words[start:end])
    return " ".join(words[:window_words*2])

def context_similarity(a: str, b: str) -> float:
    sa = {w.lower() for w in re.findall(r'\w+', a) if len(w)>2}
    sb = {w.lower() for w in re.findall(r'\w+', b) if len(w)>2}
    if not sa or not sb:
        return 0.0
    inter = sa.intersection(sb)
    union = sa.union(sb)
    return len(inter) / len(union)

def extract_number_occurrences(slides: List[Dict[str,Any]]) -> List[Dict[str,Any]]:
    occs = []
    for s in slides:
        text = s.get("text","") or ""
        numbers = s.get("numbers") or find_numbers(text)
        for tok in numbers:
            val, unit = normalize_number_token(tok)
            ctx = get_context(text, tok)
            occs.append({
                "slide": s["slide_num"],
                "raw": tok,
                "value": val,
                "unit": unit,
                "context": ctx
            })
    return occs

def detect_numeric_conflicts(slides: List[Dict[str,Any]], rel_tol: float = 0.15, ctx_thresh: float = 0.25) -> List[Dict[str,Any]]:
    occs = extract_number_occurrences(slides)
    issues = []
    seen_pairs: Set[Tuple[int,int,str]] = set()
    for i in range(len(occs)):
        a = occs[i]
        if a["value"] is not None and a["value"] < 5 and a["unit"]=='':
            continue
        for j in range(i+1, len(occs)):
            b = occs[j]
            if a["slide"] == b["slide"]:
                continue
            if a["unit"] != b["unit"]:
                continue
            if a["value"] is None or b["value"] is None:
                continue
            sim = context_similarity(a["context"], b["context"])
            if sim < ctx_thresh:
                continue
            denom = max(abs(a["value"]), abs(b["value"]), 1.0)
            rel_diff = abs(a["value"] - b["value"]) / denom
            if rel_diff > rel_tol:
                pair = tuple(sorted((a["slide"], b["slide"])))
                key = (pair[0], pair[1], a["raw"] + "::" + b["raw"])
                if key in seen_pairs:
                    continue
                seen_pairs.add(key)
                issues.append({
                    "issue_type": "numeric_mismatch",
                    "description": f"Possible numeric mismatch (ctx_sim={sim:.2f}, rel_diff={rel_diff:.2f}): Slide {a['slide']} has '{a['raw']}' (~{a['value']} {a['unit']}), Slide {b['slide']} has '{b['raw']}' (~{b['value']} {b['unit']}).",
                    "slides_involved": [a['slide'], b['slide']],
                    "source": "local_rule"
                })
    for s in slides:
        text = s.get("text","") or ""
        pct_vals = []
        for m in PERCENT_RE.finditer(text):
            try:
                pct_vals.append(float(m.group("pct")))
            except:
                pass
        if pct_vals:
            total = sum(pct_vals)
            if not (95 <= total <= 105):
                issues.append({
                    "issue_type": "percentage_sum_mismatch",
                    "description": f"Slide {s['slide_num']} percent values sum to {total:.2f}% (tolerance ±5%).",
                    "slides_involved": [s['slide_num']],
                    "source": "local_rule"
                })
    return issues

def detect_year_issues(slides: List[Dict[str,Any]]) -> List[Dict[str,Any]]:
    years_map = defaultdict(list)
    for s in slides:
        yrs = s.get("years") or find_years(s.get("text","") or "")
        for y in yrs:
            years_map[y].append(s["slide_num"])
    if len(years_map) > 1:
        return [{
            "issue_type": "multiple_years_detected",
            "description": f"Multiple years found in deck: {sorted(years_map.keys())}. Verify timeline consistency.",
            "slides_involved": sorted({sl for lst in years_map.values() for sl in lst}),
            "source": "local_rule"
        }]
    return []


def call_gemini(slides: List[Dict[str,Any]], api_key: str, max_chars_per_slide: int = 900) -> List[Dict[str,Any]]:
    try:
        genai.configure(api_key=api_key)
    except Exception as e:
        return [{
            "issue_type": "gemini_error",
            "description": f"Gemini configuration failed: {e}",
            "slides_involved": [],
            "source": "gemini"
        }]

    prompt = (
        "You are an assistant that finds factual or logical inconsistencies across presentation slides. "
        "For every issue return a JSON object with exactly these fields: "
        "\"issue_type\" (string), \"description\" (string), \"slides_involved\" (list of integers). "
        "Return a JSON ARRAY (no commentary, no markdown fences) of these issue objects.\n\nSlides:\n"
    )
    for s in slides:
        txt = (s.get("text","") or "")[:max_chars_per_slide].replace("\n"," ")
        prompt += f"Slide {s['slide_num']}: {txt}\n\n"
    prompt += "\nReturn a JSON array of issues and nothing else."

    try:
        model = genai.GenerativeModel("models/gemini-2.5-flash")
        resp = model.generate_content(prompt)
        raw_text = getattr(resp, "text", None) or getattr(resp, "output", None) or str(resp)
    except Exception as e:
        return [{
            "issue_type": "gemini_error",
            "description": f"Gemini API call failed: {e}",
            "slides_involved": [],
            "source": "gemini"
        }]

    raw = str(raw_text).strip()
    raw = re.sub(r"^```(?:json)?\s*", "", raw)
    raw = re.sub(r"```$", "", raw).strip()

    try:
        parsed = json.loads(raw)
        if isinstance(parsed, dict):
            parsed = [parsed]
        for it in parsed:
            it.setdefault("source","gemini")
        return parsed
    except Exception:
        m = re.search(r"(\[\s*\{.*\}\s*\])", raw, flags=re.DOTALL)
        if m:
            try:
                parsed = json.loads(m.group(1))
                for it in parsed:
                    it.setdefault("source","gemini")
                return parsed
            except Exception:
                pass
        return [{
            "issue_type":"gemini_error",
            "description": f"Could not parse Gemini output as JSON. Raw (truncated): {raw[:2000]}",
            "slides_involved": [],
            "source": "gemini"
        }]


def merge_issues(local: List[Dict[str,Any]], remote: List[Dict[str,Any]]) -> List[Dict[str,Any]]:
    merged = []
    seen = set()
    def sig(issue: Dict[str,Any]) -> str:
        slides = ",".join(map(str, sorted(issue.get("slides_involved",[]))))
        return f"{issue.get('issue_type')}|{slides}|{issue.get('description')[:200]}"
    for it in local + remote:
        s = sig(it)
        if s in seen:
            continue
        seen.add(s)
        merged.append(it)
    return merged

def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--pptx", "-p", required=True, help="Input .pptx file")
    parser.add_argument("--imgdir", "-d", default="slide_images", help="Directory to save extracted slide images")
    parser.add_argument("--output", "-o", default="final_report.json", help="Final issues JSON output file")
    parser.add_argument("--api_key", "-k", required=True, help="Gemini 2.5 Flash API key (mandatory)")
    parser.add_argument("--no-ocr", dest="ocr", action="store_false", help="Disable OCR pass on extracted images")
    parser.add_argument("--rel_tol", type=float, default=0.15, help="Relative tolerance for numeric mismatch")
    parser.add_argument("--ctx_thresh", type=float, default=0.25, help="Context similarity threshold")
    args = parser.parse_args()

    print(f"[+] Extracting text & images from {args.pptx} (images -> {args.imgdir}) ...")
    slides = analyze_pptx_in_memory(args.pptx, args.imgdir, ocr=args.ocr)
    print(f"[+] Extracted {len(slides)} slides. (Images saved to {args.imgdir} if present)")

    print("[+] Running local numeric and timeline checks")
    local_numeric = detect_numeric_conflicts(slides, rel_tol=args.rel_tol, ctx_thresh=args.ctx_thresh)
    local_years = detect_year_issues(slides)
    local_issues = local_numeric + local_years
    print(f"[+] Local checks found {len(local_issues)} issues")

    print("[+] Calling Gemini 2.5 Flash for semantic contradiction detection")
    gemini_issues = call_gemini(slides, api_key=args.api_key)
    print(f"[+] Gemini returned {len(gemini_issues)} items (may include gemini_error)")

    all_issues = merge_issues(local_issues, gemini_issues)
    with open(args.output, "w", encoding="utf-8") as f:
        json.dump(all_issues, f, indent=2, ensure_ascii=False)

    print(f"[+] Done. Merged issues count: {len(all_issues)}. Saved to {args.output}")
    print("[+] Short summary:")
    for idx, it in enumerate(all_issues, start=1):
        slides_str = ",".join(map(str, it.get("slides_involved", []))) or "-"
        print(f"{idx}. [{it.get('issue_type')}] Slides: {slides_str} Source: {it.get('source')}")
        print(f"   {it.get('description')[:300]}")
    print("[+] Finished.")

if __name__ == "__main__":
    main()
